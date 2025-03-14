import pandas as pd
from flask import Flask, render_template, request, send_file, abort
from pptx import Presentation
import qrcode
import requests
import os

app = Flask(__name__)

# Load Excel files
def load_data():
    try:
        people_file = pd.read_csv('People.csv')
        products_file = pd.read_csv('Products.csv')
        return people_file, products_file
    except FileNotFoundError:
        abort(404, description='CSV file not found.')

# Generate QR code image from URL
def generate_qr_code(url, output_file):
    try:
        qr = qrcode.make(url)
        qr.save(output_file)
        return output_file
    except Exception as e:
        print(f"QR code generation failed: {e}")
        return None

# Download image from URL and save as file
def download_image(url, output_file):
    try:
        response = requests.get(url)
        response.raise_for_status()
        with open(output_file, 'wb') as file:
            file.write(response.content)
        return output_file
    except requests.RequestException as e:
        print(f"Image download failed: {e}")
        return None

# Replace text, QR code, and images in PPT
def replace_text_in_ppt(people_data, products_data, template='Sample-Landscape.pptx', output_file='Updated_Presentation.pptx'):
    if not os.path.exists(template):
        abort(404, description=f'Template file {template} not found.')
    
    ppt = Presentation(template)

    qr_code_url = people_data['QRCodeURL'].values[0] if 'QRCodeURL' in people_data else None
    img_url = people_data['ImgURL'].values[0] if 'ImgURL' in people_data else None

    for slide in ppt.slides:
        for shape in slide.shapes:
            shape_name = shape.name.strip() if hasattr(shape, 'name') else ''
            if shape.has_text_frame:
                if shape_name == 'TextBox 5':
                    section_questions = [
                        products_data.get('Section 1', pd.Series([''])).values[0],
                        products_data.get('Section 2', pd.Series([''])).values[0],
                        products_data.get('Section 3', pd.Series([''])).values[0]
                    ]
                    section_answers = [
                        people_data.get('Section 1', pd.Series([''])).values[0],
                        people_data.get('Section 2', pd.Series([''])).values[0],
                        people_data.get('Section 3', pd.Series([''])).values[0]
                    ]
                    combined_text = '\n\n'.join(f"{q}:\n{a}" for q, a in zip(section_questions, section_answers))
                    shape.text_frame.clear()
                    shape.text_frame.text = combined_text

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if shape_name == 'Footer Placeholder 4':
                            run.text = str(people_data.get('FootNote', pd.Series([''])).values[0])
                        elif shape_name == 'TextBox 11':
                            run.text = str(people_data.get('Updated', pd.Series([''])).values[0])
                        elif shape_name == 'Title 1':
                            run.text = str(people_data.get('Title', pd.Series([''])).values[0])
                        elif shape_name == 'TextBox 2':
                            shape.text_frame.clear()
                            new_subtitle = str(people_data.get('SubTitle', pd.Series([''])).values[0])
                            shape.text_frame.text = new_subtitle

            if shape_name == 'QR CODE' and qr_code_url:
                img_file = generate_qr_code(qr_code_url, 'qrcode.png')
                if img_file:
                    slide.shapes.add_picture(img_file, shape.left, shape.top, shape.width, shape.height)
                    slide.shapes._spTree.remove(shape._element)

            if shape_name == 'IMAGE' and img_url:
                img_file = download_image(img_url, 'image.png')
                if img_file:
                    slide.shapes.add_picture(img_file, shape.left, shape.top, shape.width, shape.height)
                    slide.shapes._spTree.remove(shape._element)

    ppt.save(output_file)
    return output_file

@app.route('/')
def index():
    people_file, _ = load_data()
    columns = people_file.columns.tolist()
    return render_template('index.html', columns=columns)

@app.route('/general_search', methods=['POST'])
def general_search():
    search_term = request.form['general_search_term'].strip().lower()
    people_file, _ = load_data()
    
    print(f"Search term: {search_term}")
    print(f"People file preview: {people_file.head()}")

    if people_file.empty:
        return render_template('index.html', search_term=search_term, data_with_columns=[], columns=[])

    matches = people_file.apply(lambda row: row.astype(str).str.lower().str.contains(search_term).any(), axis=1)
    filtered_results = people_file[matches]

    print(f"Filtered results: {filtered_results}")

    data_with_columns = filtered_results.to_dict(orient='records')

    return render_template('index.html', search_term=search_term, data_with_columns=data_with_columns, columns=people_file.columns.tolist())

@app.route('/specific_search', methods=['POST'])
def specific_search():
    column = request.form['column']
    search_term = request.form['specific_search_term'].strip().lower()
    template = request.form['template']
    people_file, products_file = load_data()

    filtered_people = people_file[people_file[column].astype(str).str.lower() == search_term]
    if not filtered_people.empty:
        category = filtered_people.get('Category', pd.Series([''])).values[0]
        filtered_products = products_file[products_file['Category'].astype(str).str.lower() == category.lower()]
        if not filtered_products.empty:
            output_file = replace_text_in_ppt(filtered_people, filtered_products, template)
            return send_file(output_file, as_attachment=True)
        else:
            return "No matching data found in Products.csv."
    else:
        return "No matching data found in People.csv."

if __name__ == '__main__':
    app.run(debug=True)
